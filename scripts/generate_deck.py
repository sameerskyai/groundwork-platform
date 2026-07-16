#!/usr/bin/env python3
"""Generate Groundwork pitch deck (16 slides, YC style)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors
DARK_BG = RGBColor(10, 9, 8)       # #0A0908
COPPER = RGBColor(191, 122, 58)    # #BF7A3A
WHITE = RGBColor(255, 255, 255)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_slide_title_content(prs, title, content_lines, is_title_slide=False):
    """Add a slide with title and content lines (YC style: minimal text, big headline)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(60) if is_title_slide else Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Content
    if content_lines:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.4), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, line in enumerate(content_lines):
            if i > 0:
                content_frame.add_paragraph()
            para = content_frame.paragraphs[i]
            para.text = line
            para.font.size = Pt(28)
            para.font.color.rgb = WHITE
            para.space_before = Pt(12)
            para.space_after = Pt(12)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = "Sameer · Ryan · Armin · July 2026 · Confidential"
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = RGBColor(100, 100, 100)
    footer_para.alignment = PP_ALIGN.RIGHT

    return slide

# S1: Title
add_slide_title_content(prs, "GROUNDWORK",
    ["The Homeownership Layer — every home in America gets a living digital twin"],
    is_title_slide=True)

# S2: The Problem
add_slide_title_content(prs, "Finding a contractor is broken.",
    ["$500B+ market runs on guesswork",
     "No one knows who to trust or what's fair",
     "Horror stories are the norm"])

# S3: The Real Problem
add_slide_title_content(prs, "But that's not the real problem.",
    ["Renovation apps get deleted the day the project ends",
     "",
     '"A home is a 30-year relationship.',
     'Every product treats it like a one-night stand."'])

# S4: The Insight
add_slide_title_content(prs, "THE INSIGHT",
    ["Matching gets them in. MEMORY keeps them forever.",
     "",
     "Contractor matching is the front door — not the product"])

# S5: Pillar 1
add_slide_title_content(prs, "PILLAR 1: THE MATCH",
    ["AI estimate in 30 sec",
     "Budget anchored to real data",
     "5 personality questions → compatibility",
     "Swipe. Match. Message. Hire.",
     "Below 80% never shows"])

# S6: The Home Passport
add_slide_title_content(prs, "THE HOME PASSPORT",
    ["Verified permanent record of everything done to the home",
     "",
     "TRANSFERS to the buyer at sale",
     "",
     '"Carfax for homes — except we also do the repairs"',
     "6M home sales/yr = 6M free acquisition events"])

# S7: The Backstory Engine
add_slide_title_content(prs, "THE BACKSTORY ENGINE",
    ["Import the home's PAST: receipts, invoices, memory",
     "AI reads documents, camera reads appliance labels",
     "Public records auto-fill history",
     "",
     '"20 years of your home\'s history rebuilt in 10 minutes"'])

# S8: The Home Health Score
add_slide_title_content(prs, "THE HOME HEALTH SCORE",
    ["Every house gets a 0–100 score that moves with maintenance",
     "",
     "A credit score for your home",
     "Neighbors compare it. Buyers demand it. Insurers will pay for it."])

# S9: The Oracle
add_slide_title_content(prs, "THE ORACLE",
    ['AI that speaks FIRST: "Your water heater fails at year 12 in this ZIP"',
     "",
     '"Planned: $1,400. Flooded: $6,000."',
     "",
     "Every completed job makes it smarter for every home"])

# S10: The Neighborhood Nervous System
add_slide_title_content(prs, "THE NEIGHBORHOOD NERVOUS SYSTEM",
    ['"7 neighbors need driveway sealing — 30% off group rate"',
     "",
     "Storm-response triggers from real street data",
     "",
     "Homeowner demand, aggregated like a union"])

# S11: The Moat
add_slide_title_content(prs, "THE MOAT",
    ["The Match → copyable in months (the wedge)",
     "The Graph → copyable in years (needs data)",
     "Passport transfer network → NOT copyable (network effect on housing stock)",
     "Oracle accuracy → NOT copyable (our data only)"])

# S12: The Market
add_slide_title_content(prs, "THE MARKET",
    ["Old: 14M renovation events/yr (episodic)",
     "New: 145M owner-occupied homes (recurring)",
     "",
     "Path: one metro → proven retention → Passport loop compounds",
     "Ceiling: 145M × $240/yr ≈ $1.7B+ ARR — home by home, ZIP by ZIP"])

# S13: Business Model
add_slide_title_content(prs, "BUSINESS MODEL",
    ["Free: estimate, limited matching, community",
     "Homeowner+ $20/mo — $10/mo FOR LIFE at 10 referrals (CAC weapon)",
     "Contractors: free / $49 unlimited",
     "",
     "Not 'access to contractors' — the operating system for your biggest asset"])

# S14: Five Waves
add_slide_title_content(prs, "EXECUTION — FIVE WAVES",
    ["W1 NOW: The Match (~70% built, live DB, 108 tests green)",
     "W2: Graph + Passport + Backstory (wks 2–4)",
     "W3: Score + Oracle (wks 4–8)",
     "W4: Nervous System (wks 8–12)",
     "W5: Passport transfer, PM/agents, financial layer (post-raise)"])

# S15: The Playbook
add_slide_title_content(prs, "THE PLAYBOOK",
    ["Bezos: moats built on data that compounds daily",
     "Thiel: competition is for losers — build the monopoly layer",
     "Buffett: own assets that appreciate — the Graph appreciates with every job",
     "YC: make something people want, then make leaving unthinkable"])

# S16: Final Slide
add_slide_title_content(prs, "WE'RE BUILDING THE LAYER THE AMERICAN HOME RUNS ON.",
    ["The Match gets them in",
     "The Graph makes leaving unthinkable",
     "The Oracle makes us indispensable",
     "The Passport makes every home sale our salesman",
     "",
     "Ship Wave 1. Everything follows."])

# Save
output_path = "/Users/sameerhersi/Desktop/Groundwork_Pitch_Deck.pptx"
prs.save(output_path)
print(f"✅ Deck saved to {output_path}")
